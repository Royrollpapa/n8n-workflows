import os
from pathlib import Path
import re
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def extract_numeric_blocks(md_content: str):
    results = []
    unsupported_blocks = []
    # 1. 识别标准Markdown表格
    table_pattern = re.compile(r'(\|.+\|\n(?:\|[-: ]+\|\n)+((?:\|.*\|\n?)+))', re.MULTILINE)
    for match in table_pattern.finditer(md_content):
        raw = match.group(0)
        try:
            lines = [line.strip() for line in raw.strip().split('\n') if line.strip()]
            # 跳过分隔线
            data_lines = [line for line in lines if not re.match(r'^\|?[-: ]+\|?$', line)]
            if len(data_lines) < 2:
                unsupported_blocks.append(raw)
                continue
            header = [cell.strip() for cell in data_lines[0].strip('|').split('|')]
            rows = []
            for line in data_lines[1:]:
                row = [cell.strip() for cell in line.strip('|').split('|')]
                if len(row) == len(header):
                    rows.append(row)
            df = pd.DataFrame(rows, columns=header)
            # 尝试将数值列转为float
            for col in df.columns:
                try:
                    df[col] = df[col].astype(float)
                except Exception:
                    pass
            results.append({'type': 'table', 'data': df, 'raw': raw})
        except Exception:
            unsupported_blocks.append(raw)
            continue
    # 2. 识别有数字的无序列表
    list_pattern = re.compile(r'(?:^|\n)(- .*[0-9]+.*(?:\n- .*[0-9]+.*)+)', re.MULTILINE)
    for match in list_pattern.finditer(md_content):
        raw = match.group(1)
        lines = [line.strip('- ').strip() for line in raw.strip().split('\n') if line.strip()]
        results.append({'type': 'list', 'data': lines, 'raw': raw})
    # 3. 识别代码块中的数字数据
    code_pattern = re.compile(r'```[a-zA-Z0-9]*\n([\s\S]*?\d+[\s\S]*?)```', re.MULTILINE)
    for match in code_pattern.finditer(md_content):
        raw = match.group(1)
        lines = [line for line in raw.strip().split('\n') if any(c.isdigit() for c in line)]
        if lines:
            results.append({'type': 'code', 'data': lines, 'raw': raw})
    # 检查未被识别的数字块
    if unsupported_blocks:
        for block in unsupported_blocks:
            print(f"[提示] 检测到无法识别的数字表格或数据块，请检查格式：\n{block}\n建议参考README中的数字数据格式示例进行调整。\n")
    if not results:
        print("[提示] 未检测到可识别的数字数据块。请参考README中的数字数据格式示例进行编写。")
    return results

def auto_chart_type(block):
    if block['type'] == 'table':
        df = block['data']
        if df.shape[1] == 2:
            col1, col2 = df.columns[0], df.columns[1]
            if all(isinstance(x, (int, float, np.integer, np.floating)) or str(x).replace('.','',1).isdigit() for x in df[col2]):
                if any(s in str(col1) for s in ['年', '月', 'date', 'time']) or all(str(x).isdigit() for x in df[col1]):
                    return 'line'
                else:
                    return 'bar'
        elif df.shape[1] > 2:
            if any('占比' in str(c) or '比例' in str(c) for c in df.columns):
                return 'pie'
            else:
                return 'stacked_bar'
    if block['type'] == 'list':
        return 'bar'
    if block['type'] == 'code':
        lines = block['data']
        if all(',' in line for line in lines):
            return 'line'
        else:
            return 'bar'
    return 'table'

def render_auto_chart(block):
    chart_type = auto_chart_type(block)
    fig, ax = plt.subplots(figsize=(8, 5))
    img_bytes = None
    try:
        if chart_type == 'bar':
            if block['type'] == 'table':
                df = block['data']
                ax.bar(df.iloc[:, 0], df.iloc[:, 1])
                ax.set_title(f"{df.columns[0]} vs {df.columns[1]}")
            elif block['type'] == 'list':
                labels = [str(i+1) for i in range(len(block['data']))]
                values = [float(''.join(filter(str.isdigit, s))) for s in block['data']]
                ax.bar(labels, values)
                ax.set_title("数据分布")
            elif block['type'] == 'code':
                lines = [line.split(',') for line in block['data'] if ',' in line]
                labels = [l[0] for l in lines]
                values = [float(l[1]) for l in lines]
                ax.bar(labels, values)
                ax.set_title("数据分布")
        elif chart_type == 'line':
            if block['type'] == 'table':
                df = block['data']
                ax.plot(df.iloc[:, 0], df.iloc[:, 1], marker='o')
                ax.set_title(f"{df.columns[0]} vs {df.columns[1]}")
            elif block['type'] == 'code':
                lines = [line.split(',') for line in block['data'] if ',' in line]
                labels = [l[0] for l in lines]
                values = [float(l[1]) for l in lines]
                ax.plot(labels, values, marker='o')
                ax.set_title("趋势分析")
        elif chart_type == 'stacked_bar':
            df = block['data']
            x = df.iloc[:, 0]
            y = df.iloc[:, 1:]
            y.cumsum(axis=1).plot(kind='bar', stacked=True, ax=ax)
            ax.set_xticklabels(x)
            ax.set_title("多组数据对比")
        elif chart_type == 'pie':
            df = block['data']
            ax.pie(df.iloc[:, 1], labels=df.iloc[:, 0], autopct='%1.1f%%')
            ax.set_title("占比分析")
        else:
            plt.close(fig)
            return None
        buf = io.BytesIO()
        plt.tight_layout()
        plt.savefig(buf, format='png')
        plt.close(fig)
        img_bytes = buf.getvalue()
    except Exception as e:
        plt.close(fig)
        return None
    return img_bytes

def add_table_to_slide(slide, block, left=Inches(0.5), top=Inches(4.5)):
    if block['type'] == 'table':
        df = block['data']
        rows, cols = df.shape
        table = slide.shapes.add_table(rows+1, cols, left, top, Inches(8), Inches(1.2)).table
        for j, col in enumerate(df.columns):
            table.cell(0, j).text = str(col)
        for i in range(rows):
            for j in range(cols):
                table.cell(i+1, j).text = str(df.iloc[i, j])
    elif block['type'] == 'list':
        lines = block['data']
        table = slide.shapes.add_table(len(lines)+1, 2, left, top, Inches(8), Inches(1.2)).table
        table.cell(0, 0).text = '序号'
        table.cell(0, 1).text = '内容'
        for i, item in enumerate(lines):
            table.cell(i+1, 0).text = str(i+1)
            table.cell(i+1, 1).text = str(item)
    elif block['type'] == 'code':
        lines = block['data']
        table = slide.shapes.add_table(len(lines)+1, 1, left, top, Inches(8), Inches(1.2)).table
        table.cell(0, 0).text = '内容'
        for i, line in enumerate(lines):
            table.cell(i+1, 0).text = str(line)

def md_to_ppt(md_path, output_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        md_content = f.read()
    numeric_blocks = extract_numeric_blocks(md_content)
    prs = Presentation()
    for idx, block in enumerate(numeric_blocks):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # 标题
        title_shape = slide.shapes.title
        title_shape.text = f"自动图表与表格 - {idx+1}"
        # 图表
        img_bytes = render_auto_chart(block)
        if img_bytes:
            pic = slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(1.2), Inches(6), Inches(3))
        # 表格
        add_table_to_slide(slide, block)
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")

if __name__ == "__main__":
    input_dir = Path(__file__).parent.parent / 'input'
    output_dir = Path(__file__).parent.parent / 'output'
    for file in os.listdir(input_dir):
        if file.endswith('.md'):
            md_path = input_dir / file
            pptx_path = output_dir / (Path(file).stem + '.pptx')
            md_to_ppt(str(md_path), str(pptx_path)) 