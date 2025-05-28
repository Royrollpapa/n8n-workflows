"""
演示文稿处理器
处理PowerPoint文档的创建和转换
"""

import os
from pathlib import Path
from typing import Optional, Dict, Any
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PresentationProcessor:
    """演示文稿处理器类"""
    
    def __init__(self):
        self.default_font = 'Microsoft YaHei'
        self.title_size = Pt(44)
        self.subtitle_size = Pt(32)
        self.content_size = Pt(24)
        
    def from_html(self, input_path: str, output_path: str,
                  template_path: Optional[str] = None,
                  **kwargs) -> str:
        """
        从HTML创建PowerPoint演示文稿
        
        Args:
            input_path: 输入HTML文件路径
            output_path: 输出PPT文件路径
            template_path: PPT模板路径（可选）
            **kwargs: 其他参数
            
        Returns:
            str: 输出文件路径
        """
        # 创建或加载演示文稿
        prs = Presentation(template_path) if template_path else Presentation()
        
        # 读取HTML文件
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
            
        # 解析HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 处理每个主要部分
        for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'table']):
            if element.name.startswith('h'):
                # 创建新幻灯片
                slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题
                title = slide.shapes.title
                title.text = element.get_text()
                title.text_frame.paragraphs[0].font.name = self.default_font
                title.text_frame.paragraphs[0].font.size = self.title_size
                
                # 获取内容占位符
                content = slide.placeholders[1]
                
            elif element.name == 'p':
                if 'content' in locals():
                    p = content.text_frame.add_paragraph()
                    p.text = element.get_text()
                    p.font.name = self.default_font
                    p.font.size = self.content_size
                    
            elif element.name in ['ul', 'ol']:
                if 'content' in locals():
                    for li in element.find_all('li'):
                        p = content.text_frame.add_paragraph()
                        p.text = f"• {li.get_text()}" if element.name == 'ul' else f"1. {li.get_text()}"
                        p.font.name = self.default_font
                        p.font.size = self.content_size
                        
            elif element.name == 'table':
                if 'content' in locals():
                    # 创建表格
                    rows = len(element.find_all('tr'))
                    cols = len(element.find_all('th'))
                    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
                    
                    # 设置表头
                    for i, th in enumerate(element.find_all('th')):
                        cell = table.cell(0, i)
                        cell.text = th.get_text()
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.name = self.default_font
                        cell.text_frame.paragraphs[0].font.size = self.content_size
                    
                    # 设置数据行
                    for i, tr in enumerate(element.find_all('tr')[1:], 1):
                        for j, td in enumerate(tr.find_all('td')):
                            cell = table.cell(i, j)
                            cell.text = td.get_text()
                            cell.text_frame.paragraphs[0].font.name = self.default_font
                            cell.text_frame.paragraphs[0].font.size = self.content_size
        
        # 保存演示文稿
        prs.save(output_path)
        return output_path
    
    def apply_theme(self, prs: Presentation, theme: Dict[str, Any]):
        """
        应用主题到演示文稿
        
        Args:
            prs: Presentation对象
            theme: 主题配置字典
        """
        # 设置默认字体
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = theme.get('font', self.default_font)
                        if paragraph.text == shape.text_frame.paragraphs[0].text:  # 标题
                            paragraph.font.size = theme.get('title_size', self.title_size)
                        else:  # 内容
                            paragraph.font.size = theme.get('content_size', self.content_size)
                            
        return prs 