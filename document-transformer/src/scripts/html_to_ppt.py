import os
import yaml
import logging
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class HTMLToPPT:
    """HTML转PPT转换器"""
    
    def __init__(self, config_path: str):
        """
        初始化转换器
        Args:
            config_path: 配置文件路径
        """
        self.config_path = str(Path(__file__).parent / ".." / ".." / "config" / "config.yaml")
        with open(self.config_path, 'r', encoding='utf-8') as f:
            self.config = yaml.safe_load(f)
        
    def _check_file_permissions(self, file_path):
        """检查文件权限"""
        try:
            # 检查目录是否存在
            directory = os.path.dirname(file_path)
            if not os.path.exists(directory):
                os.makedirs(directory)
            
            # 检查文件是否可写
            if os.path.exists(file_path):
                if not os.access(file_path, os.W_OK):
                    raise PermissionError(f"文件 {file_path} 没有写入权限")
            else:
                # 测试目录是否可写
                test_file = os.path.join(directory, 'test.txt')
                with open(test_file, 'w') as f:
                    f.write('test')
                os.remove(test_file)
                
        except Exception as e:
            logger.error(f"文件权限检查失败: {str(e)}")
            raise

    def convert_to_ppt(self, html_path: str, output_path: str) -> str:
        """
        将HTML文件转换为PPT
        Args:
            html_path: HTML文件路径
            output_path: 输出文件路径
        Returns:
            str: 输出文件路径
        """
        try:
            # 检查文件权限
            self._check_file_permissions(html_path)
            self._check_file_permissions(output_path)
            
            # 读取HTML文件
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 创建PPT
            prs = Presentation()
            
            # 设置幻灯片大小
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # 处理每个幻灯片内容
            for slide_content in soup.find_all('div', class_='slide-content'):
                # 创建新幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # 设置标题
                title = slide_content.find('h1')
                if title:
                    title_shape = slide.shapes.title
                    title_shape.text = title.text
                    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
                    title_shape.text_frame.paragraphs[0].font.bold = True
                    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 设置内容
                # 创建文本框
                left = Inches(1)
                top = Inches(2)
                width = Inches(14)
                height = Inches(6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                
                # 添加段落
                for p in slide_content.find_all(['p', 'h2', 'h3', 'ul', 'ol']):
                    if p.name == 'p':
                        paragraph = tf.add_paragraph()
                        paragraph.text = p.text
                        paragraph.font.size = Pt(24)
                    elif p.name == 'h2':
                        paragraph = tf.add_paragraph()
                        paragraph.text = p.text
                        paragraph.font.size = Pt(32)
                        paragraph.font.bold = True
                    elif p.name == 'h3':
                        paragraph = tf.add_paragraph()
                        paragraph.text = p.text
                        paragraph.font.size = Pt(28)
                        paragraph.font.bold = True
                    elif p.name in ['ul', 'ol']:
                        for li in p.find_all('li'):
                            paragraph = tf.add_paragraph()
                            paragraph.text = f"• {li.text}"
                            paragraph.font.size = Pt(24)
            
            # 保存PPT
            prs.save(output_path)
            logger.info(f"PPT生成成功: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}")
            raise 